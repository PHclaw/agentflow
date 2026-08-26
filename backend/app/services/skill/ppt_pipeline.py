"""PPT Generation Skill：根据用户大纲在服务端生成可下载 PPTX。"""
from __future__ import annotations

import json
import re
import uuid
from pathlib import Path
from typing import Any

from app.logging_setup import get_logger

logger = get_logger("ppt")

from app.core.paths import generated_root

GENERATED_DIR = generated_root()
PPT_FONT = "Microsoft YaHei"


def _set_run_font(paragraph, *, name: str = PPT_FONT) -> None:
    from lxml import etree
    from pptx.oxml.ns import qn

    for run in paragraph.runs:
        run.font.name = name
        r_pr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            node = r_pr.find(qn(tag))
            if node is None:
                node = etree.SubElement(r_pr, qn(tag))
            node.set("typeface", name)

STYLES = {
    "keynote",
    "dark-premium",
    "minimal-swiss",
    "glassmorphism",
    "gradient-modern",
    "editorial",
    "business",
    "academic",
    "minimal",
    "creative",
}


def is_ppt_skill(agent) -> bool:
    wf = getattr(agent, "workflow", None) or {}
    if isinstance(wf, dict) and wf.get("kind") in {"ppt", "ppt-generation"}:
        return True
    sid = (getattr(agent, "id", None) or "").lower()
    if sid in {"ppt-generation", "ppt"}:
        return True
    specialty = getattr(agent, "specialty", None) or ""
    name = getattr(agent, "name", None) or ""
    blob = f"{specialty} {name}".lower()
    return any(k in blob for k in ("ppt", "演示文稿", "幻灯片", "powerpoint"))


def _public_url(filename: str) -> str:
    return f"/static/generated/{filename}"


def _prefer_current(text: str) -> str:
    t = text or ""
    marker = "## 当前用户"
    if marker in t:
        return t.split(marker, 1)[-1].strip()
    return t.strip()


def _detect_style(text: str) -> str:
    t = (text or "").lower()
    for s in (
        "glassmorphism",
        "dark-premium",
        "gradient-modern",
        "neo-brutalist",
        "3d-isometric",
        "editorial",
        "minimal-swiss",
        "keynote",
        "academic",
        "business",
        "minimal",
        "creative",
    ):
        if s in t:
            return s
    if "学术" in (text or ""):
        return "academic"
    if "苹果" in (text or "") or "keynote" in t:
        return "keynote"
    if "极简" in (text or ""):
        return "minimal-swiss"
    if "商务" in (text or "") or "汇报" in (text or "") or "科研" in (text or ""):
        return "business"
    return "business"


def _try_parse_json_plan(text: str) -> dict[str, Any] | None:
    t = text or ""
    # fenced json
    m = re.search(r"```(?:json)?\s*(\{[\s\S]*?\})\s*```", t)
    raw = m.group(1) if m else None
    if not raw:
        # bare object with slides
        m2 = re.search(r"(\{\s*\"title\"[\s\S]*\"slides\"\s*:\s*\[[\s\S]*\]\s*\})", t)
        raw = m2.group(1) if m2 else None
    if not raw:
        return None
    try:
        plan = json.loads(raw)
    except json.JSONDecodeError:
        return None
    if isinstance(plan, dict) and isinstance(plan.get("slides"), list) and plan["slides"]:
        return plan
    return None


def _parse_markdown_outline(text: str) -> dict[str, Any] | None:
    lines = [ln.rstrip() for ln in (text or "").splitlines()]
    title = ""
    style = _detect_style(text)
    slides: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None

    title_m = re.search(r"(?:标题|题目|主题|title)\s*[：:=]\s*(.+)", text or "", flags=re.I)
    if title_m:
        title = title_m.group(1).strip().strip("「」\"'")

    for ln in lines:
        s = ln.strip()
        if not s:
            continue
        if re.match(r"^(请|帮我|生成|做|制作|创建).*(ppt|pptx|演示|幻灯片)", s, flags=re.I):
            # try extract topic after 关于/：
            tm = re.search(r"(?:关于|：|:)\s*(.+)$", s)
            if tm and not title:
                title = re.sub(r"(的)?(PPT|pptx|演示文稿|幻灯片).*$", "", tm.group(1), flags=re.I).strip()
            continue
        if re.match(r"^(风格|style)\s*[：:=]", s, flags=re.I):
            continue
        heading = re.match(r"^(?:#{1,3}\s+|第\s*\d+\s*[页章节部分]\s*[：:.]?\s*|\d+[\.、]\s*)(.+)$", s)
        if heading:
            if current:
                slides.append(current)
            h = heading.group(1).strip()
            current = {"type": "content", "title": h, "key_points": []}
            continue
        bullet = re.match(r"^[-*•·]\s+(.+)$", s) or re.match(r"^[（(]?\d+[）)]\s*(.+)$", s)
        if bullet and current is not None:
            current.setdefault("key_points", []).append(bullet.group(1).strip())
            continue
        if current is not None and len(s) <= 80 and not s.endswith("。"):
            # treat short line as bullet
            current.setdefault("key_points", []).append(s)

    if current:
        slides.append(current)

    if not slides:
        return None
    if not title:
        title = str(slides[0].get("title") or "演示文稿")
        if slides[0].get("type") != "title":
            slides[0]["type"] = "title"
            if not slides[0].get("subtitle") and slides[0].get("key_points"):
                slides[0]["subtitle"] = " · ".join(slides[0]["key_points"][:2])
                slides[0]["key_points"] = []
    else:
        # ensure first is title slide if missing
        if slides and slides[0].get("type") != "title":
            slides.insert(0, {"type": "title", "title": title, "subtitle": "", "key_points": []})

    for i, sl in enumerate(slides, 1):
        sl["slide_number"] = i
        sl.setdefault("type", "title" if i == 1 else "content")
        sl.setdefault("key_points", [])
    return {
        "title": title,
        "style": style,
        "aspect_ratio": "16:9",
        "slides": slides[:20],
    }


def _fallback_plan(text: str) -> dict[str, Any]:
    style = _detect_style(text)
    topic = ""
    m = re.search(
        r"(?:关于|主题|标题|做[一]?个|生成|制作)\s*[：:]?\s*(.+?)(?:的)?(?:PPT|pptx|演示|幻灯片|$)",
        text or "",
        flags=re.I,
    )
    if m:
        topic = m.group(1).strip(" 「」\"'。，,")
    if not topic:
        # take first non-empty short line
        for ln in (text or "").splitlines():
            s = ln.strip()
            if s and len(s) <= 40 and "生成" not in s:
                topic = s
                break
    topic = topic or "演示文稿"
    return {
        "title": topic,
        "style": style,
        "aspect_ratio": "16:9",
        "slides": [
            {
                "slide_number": 1,
                "type": "title",
                "title": topic,
                "subtitle": "由平台 PPT Skill 生成",
                "key_points": [],
            },
            {
                "slide_number": 2,
                "type": "content",
                "title": "议程",
                "key_points": ["背景与目标", "核心方案", "下一步行动"],
            },
            {
                "slide_number": 3,
                "type": "content",
                "title": "核心要点",
                "key_points": ["要点一：请补充具体内容", "要点二：请补充具体内容", "要点三：请补充具体内容"],
            },
            {
                "slide_number": 4,
                "type": "conclusion",
                "title": "总结与行动",
                "subtitle": "Thank you",
                "key_points": ["确认负责人", "明确时间节点", "收集反馈"],
            },
        ],
        "inferred": True,
    }


def build_plan(user_text: str) -> dict[str, Any]:
    text = _prefer_current(user_text)
    plan = _try_parse_json_plan(text)
    if plan:
        plan.setdefault("style", _detect_style(text))
        plan.setdefault("aspect_ratio", "16:9")
        return plan
    plan = _parse_markdown_outline(text)
    if plan:
        return plan
    return _fallback_plan(text)


def _theme(style: str) -> dict[str, Any]:
    s = (style or "corporate-clean").lower()
    if s in {
        "dark-premium",
        "keynote",
        "glassmorphism",
        "gradient-modern",
        "creative",
        "tokyo-night",
        "catppuccin-mocha",
        "dracula",
        "cyberpunk-neon",
    }:
        return {
            "bg": (10, 10, 12),
            "title": (255, 255, 255),
            "body": (210, 214, 220),
            "accent": (0, 113, 227) if s == "keynote" else (0, 212, 255),
        }
    if s in {"academic", "editorial", "editorial-serif", "academic-paper"}:
        return {
            "bg": (250, 249, 246),
            "title": (32, 32, 36),
            "body": (55, 55, 60),
            "accent": (122, 45, 18),
        }
    # corporate-clean / swiss / 蓝灰商务
    return {
        "bg": (255, 255, 255),
        "title": (10, 37, 64),
        "body": (66, 84, 102),
        "accent": (29, 78, 216),
    }


def generate_pptx(plan: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation

    from app.services.skill.ppt_office import build_editable_pptx

    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    slides = plan.get("slides") or []
    if not slides:
        return {
            "intent": "generate",
            "exitCode": 1,
            "note": "计划中没有 slides。",
            "stdout": "",
            "stderr": "",
        }

    prs = Presentation()
    build_editable_pptx(plan, prs)

    fname = f"ppt-{uuid.uuid4().hex[:10]}.pptx"
    out_path = GENERATED_DIR / fname
    prs.save(str(out_path))
    url = _public_url(fname)
    stdout = (
        f"title={plan.get('title')}\n"
        f"style={plan.get('style')}\n"
        f"slides={len(slides)}\n"
        f"inferred={bool(plan.get('inferred'))}\n"
        f"downloadUrl={url}\n"
        f"planJson={json.dumps({'title': plan.get('title'), 'style': plan.get('style'), 'slideTitles': [s.get('title') for s in slides]}, ensure_ascii=False)}"
    )
    return {
        "intent": "generate",
        "script": "python-pptx",
        "exitCode": 0,
        "stdout": stdout,
        "stderr": "",
        "note": "已生成可编辑 PPTX（原生文本/表格/簇状柱状图，非截图）",
        "downloadUrl": url,
        "downloadName": fname,
        "slideCount": len(slides),
        "title": plan.get("title"),
        "style": plan.get("style"),
        "inferred": bool(plan.get("inferred")),
    }


_IMAGE_SUFFIX = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}


def _pptx_files(uploaded: list[dict] | None) -> list[dict]:
    out: list[dict] = []
    for f in uploaded or []:
        name = str(f.get("name") or "")
        path = str(f.get("path") or "")
        if name.lower().endswith(".pptx") and path and Path(path).is_file():
            out.append(f)
    return out


def _image_files(uploaded: list[dict] | None) -> list[dict]:
    out: list[dict] = []
    for f in uploaded or []:
        name = str(f.get("name") or "")
        path = str(f.get("path") or "")
        if Path(name).suffix.lower() in _IMAGE_SUFFIX and path and Path(path).is_file():
            out.append(f)
    return out


def _is_followup(text: str) -> bool:
    current = _prefer_current(text)
    if any(k in current for k in ("上文", "上面", "上述", "刚才", "刚刚", "之前的", "先前")):
        return True
    return "## 先前对话" in (text or "")


def _wants_read(text: str) -> bool:
    t = _prefer_current(text)
    return any(k in t for k in ("读", "大纲", "有几页", "讲了什么", "内容是", "outline", "有哪些页"))


def _parse_replace_pair(text: str) -> tuple[str, str] | None:
    t = _prefer_current(text)
    patterns = (
        r"把[「\"“']?(.+?)[」\"”']?(?:改成|换成|替换为)[「\"“']?(.+?)[」\"”']?\s*$",
        r"(?:替换|replace)\s+[「\"“']?(.+?)[」\"”']?\s+(?:为|成|with|to)\s+[「\"“']?(.+?)[」\"”']?",
    )
    for pat in patterns:
        m = re.search(pat, t, flags=re.I)
        if m:
            old, new = m.group(1).strip(), m.group(2).strip()
            if old and new and old != new:
                return old, new
    return None


def _parse_template_values(text: str) -> dict[str, str] | None:
    t = text or ""
    m = re.search(r"```(?:json)?\s*(\{[\s\S]*?\})\s*```", t)
    raw = m.group(1) if m else None
    if not raw:
        m2 = re.search(r"(\{\s*\"[^\"]+\"\s*:\s*\"[\s\S]*?\})", t)
        raw = m2.group(1) if m2 else None
    if not raw:
        return None
    try:
        obj = json.loads(raw)
    except json.JSONDecodeError:
        return None
    if not isinstance(obj, dict) or obj.get("slides"):
        return None
    out = {str(k): str(v) for k, v in obj.items() if not isinstance(v, (dict, list))}
    return out or None


def _save_prs(prs, prefix: str) -> tuple[str, str]:
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    fname = f"{prefix}-{uuid.uuid4().hex[:10]}.pptx"
    out_path = GENERATED_DIR / fname
    prs.save(str(out_path))
    return fname, _public_url(fname)


def _iter_text_frames(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame
        if getattr(shape, "has_table", False):
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    yield cell.text_frame
    if getattr(slide, "has_notes_slide", False):
        notes = slide.notes_slide.notes_text_frame
        if notes is not None:
            yield notes


def read_pptx(file: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation

    path = Path(str(file.get("path") or ""))
    prs = Presentation(str(path))
    slides: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = (shape.text_frame.text or "").strip()
                if t:
                    texts.append(t)
        notes = ""
        if getattr(slide, "has_notes_slide", False):
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        slides.append({"index": i, "texts": texts[:12], "notes": notes[:500]})
    outline = {"file": file.get("name"), "slideCount": len(slides), "slides": slides}
    return {
        "intent": "read",
        "script": "python-pptx-read",
        "exitCode": 0,
        "stdout": json.dumps(outline, ensure_ascii=False, indent=2)[:8000],
        "stderr": "",
        "note": f"已读取 {len(slides)} 页大纲（未改文件）",
        "slideCount": len(slides),
        "title": file.get("name"),
    }


def edit_replace_pptx(file: dict[str, Any], old: str, new: str) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(file.get("path")))
    n = 0
    for slide in prs.slides:
        for tf in _iter_text_frames(slide):
            for p in tf.paragraphs:
                for run in p.runs:
                    if old in (run.text or ""):
                        run.text = (run.text or "").replace(old, new)
                        n += 1
    fname, url = _save_prs(prs, "ppt-edit")
    return {
        "intent": "edit",
        "script": "python-pptx-edit",
        "exitCode": 0,
        "stdout": f"replaced={old!r} -> {new!r}\ncount={n}\ndownloadUrl={url}",
        "stderr": "",
        "note": f"已替换 {n} 处并生成新 PPTX",
        "downloadUrl": url,
        "downloadName": fname,
    }


def fill_template_pptx(file: dict[str, Any], values: dict[str, str]) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(file.get("path")))
    n = 0
    for slide in prs.slides:
        for tf in _iter_text_frames(slide):
            for p in tf.paragraphs:
                for run in p.runs:
                    text = run.text or ""
                    nxt = text
                    for k, v in values.items():
                        nxt = nxt.replace("{{" + k + "}}", v)
                        nxt = nxt.replace("{" + k + "}", v)
                    if nxt != text:
                        run.text = nxt
                        n += 1
    fname, url = _save_prs(prs, "ppt-tpl")
    return {
        "intent": "template",
        "script": "python-pptx-template",
        "exitCode": 0,
        "stdout": f"tokens={list(values)}\nfilledRuns={n}\ndownloadUrl={url}",
        "stderr": "",
        "note": "已按占位符填充模板并生成新 PPTX",
        "downloadUrl": url,
        "downloadName": fname,
    }


def generate_pptx_from_images(plan: dict[str, Any], images: list[dict]) -> dict[str, Any]:
    from io import BytesIO

    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    aspect = plan.get("aspect_ratio") or "16:9"
    if aspect == "4:3":
        width, height = Inches(10), Inches(7.5)
    else:
        width, height = Inches(13.333), Inches(7.5)
    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height
    blank = prs.slide_layouts[6]
    slides_info = plan.get("slides") or []
    for i, imgf in enumerate(images[:20]):
        slide = prs.slides.add_slide(blank)
        path = str(imgf.get("path"))
        with Image.open(path) as img:
            if img.mode in {"RGBA", "P"}:
                img = img.convert("RGB")
            buf = BytesIO()
            img.save(buf, format="JPEG", quality=92)
            buf.seek(0)
            slide.shapes.add_picture(buf, Inches(0), Inches(0), width=width, height=height)
        if i < len(slides_info):
            info = slides_info[i]
            notes = [str(info.get("title") or "")]
            notes.extend(str(p) for p in (info.get("key_points") or [])[:6])
            body = "\n".join(x for x in notes if x).strip()
            if body:
                slide.notes_slide.notes_text_frame.text = body
    fname, url = _save_prs(prs, "ppt-img")
    return {
        "intent": "generate-images",
        "script": "python-pptx-images",
        "exitCode": 0,
        "stdout": f"slides={len(images)}\ndownloadUrl={url}",
        "stderr": "",
        "note": "已用上传图片合成 PPTX",
        "downloadUrl": url,
        "downloadName": fname,
        "slideCount": min(len(images), 20),
        "title": plan.get("title"),
        "style": plan.get("style"),
    }


def run_ppt_tools(
    user_text: str,
    *,
    uploaded_files: list[dict] | None = None,
    conversation_id: str | None = None,
) -> dict[str, Any]:
    try:
        from app.services.skill.ppt_wizard import run_wizard

        pptx = _pptx_files(uploaded_files)
        images = _image_files(uploaded_files)
        current = _prefer_current(user_text)

        if pptx:
            pair = _parse_replace_pair(current)
            if pair:
                return edit_replace_pptx(pptx[0], pair[0], pair[1])
            values = _parse_template_values(current)
            if values:
                return fill_template_pptx(pptx[0], values)
            if _wants_read(current) or (
                not _try_parse_json_plan(current) and not _parse_markdown_outline(current)
            ):
                return read_pptx(pptx[0])
        if images and not pptx:
            plan = build_plan(user_text)
            return generate_pptx_from_images(plan, images)
        return run_wizard(
            user_text,
            conversation_id=conversation_id or "",
            plan_builder=build_plan,
            pptx_builder=generate_pptx,
        )
    except Exception as exc:  # noqa: BLE001
        logger.error("ppt tool fail: %s", exc)
        return {
            "intent": "generate",
            "exitCode": 1,
            "script": "ppt_pipeline",
            "stdout": "",
            "stderr": str(exc),
            "note": f"PPT 处理失败: {exc}",
        }


def format_ppt_tool_result(trace: dict[str, Any]) -> str:
    if not trace:
        return "（无脚本结果）"
    parts = [
        f"intent: {trace.get('intent')}",
        f"script: {trace.get('script')}",
        f"exitCode: {trace.get('exitCode')}",
    ]
    if trace.get("note"):
        parts.append(f"note: {trace['note']}")
    if trace.get("downloadUrl"):
        parts.append(f"downloadUrl: {trace['downloadUrl']}")
    if trace.get("downloadName"):
        parts.append(f"downloadName: {trace['downloadName']}")
    if trace.get("downloadUrls"):
        parts.append("downloadUrls:\n" + "\n".join(str(u) for u in trace["downloadUrls"]))
    if trace.get("stdout"):
        parts.append("stdout:\n" + str(trace["stdout"])[:8000])
    if trace.get("stderr"):
        parts.append("stderr:\n" + str(trace["stderr"])[:2000])
    parts.append(
        "重要：intent=brief / content-ask：保留脚本中的题号 1/2/3/4 与 a/b/c/d，提示用户用 `1b 2a 3a 4b` 这种方式输入，不要改成点选按钮，也不要改成开放问答题。禁止现在生成完整 PPT。"
        "intent=style-preview：对比三套风格，给出 HTML 链接，提示输入 `1a`/`1b`/`1c`（或 style-a/b/c）。"
        "intent=compose 不会展示给用户。"
        "intent=generate 且含 downloadUrl：必须给出 HTML 演示稿链接，可另附 PPTX。"
        "intent=read 只报告大纲；intent=intro 只介绍能力。"
        "禁止把系统时钟整段复制进回复。"
    )
    return "\n".join(parts)
