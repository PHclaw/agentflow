"""可编辑 Office PPTX：把 html-ppt / weekly-report 版式落到原生形状、表格、图表。"""
from __future__ import annotations

import re
from typing import Any

PPT_FONT = "Microsoft YaHei"


def _set_run_font(paragraph, *, name: str = PPT_FONT) -> None:
    from lxml import etree
    from pptx.oxml.ns import qn

    for run in paragraph.runs:
        run.font.name = name
        r_pr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            node = r_pr.find(qn(tag))
            if node is None:
                node = etree.SubElement(r_pr, qn(tag))
            node.set("typeface", name)


def tokens(style: str | None) -> dict[str, Any]:
    s = (style or "corporate-clean").lower()
    if s in {
        "tokyo-night",
        "catppuccin-mocha",
        "dracula",
        "cyberpunk-neon",
        "keynote",
        "dark-premium",
        "glassmorphism",
    }:
        return {
            "bg": (13, 17, 23),
            "soft": (22, 27, 34),
            "surface": (22, 27, 34),
            "title": (230, 237, 243),
            "body": (139, 148, 158),
            "muted": (110, 118, 129),
            "accent": (46, 160, 67),
            "accent2": (121, 192, 255),
            "line": (48, 54, 61),
            "header_fg": (13, 17, 23),
            "good": (126, 231, 135),
            "series": [(126, 231, 135), (121, 192, 255), (210, 168, 255), (255, 123, 114)],
        }
    if s in {"academic", "academic-paper", "editorial", "editorial-serif"}:
        return {
            "bg": (250, 249, 246),
            "soft": (244, 241, 234),
            "surface": (255, 255, 255),
            "title": (32, 32, 36),
            "body": (55, 55, 60),
            "muted": (122, 114, 104),
            "accent": (122, 45, 18),
            "accent2": (154, 68, 32),
            "line": (220, 214, 204),
            "header_fg": (255, 255, 255),
            "good": (26, 127, 84),
            "series": [(122, 45, 18), (90, 90, 96), (154, 68, 32), (60, 90, 120)],
        }
    # corporate-clean / weekly-report 蓝灰商务科研
    return {
        "bg": (250, 251, 252),
        "soft": (243, 245, 249),
        "surface": (255, 255, 255),
        "title": (22, 30, 55),
        "body": (80, 88, 107),
        "muted": (139, 146, 165),
        "accent": (46, 99, 235),
        "accent2": (14, 165, 181),
        "line": (226, 230, 237),
        "header_fg": (255, 255, 255),
        "good": (16, 185, 129),
        "series": [(46, 99, 235), (14, 165, 181), (245, 158, 11), (16, 185, 129)],
    }


def _rgb(c):
    from pptx.dml.color import RGBColor

    return RGBColor(*c)


def _solid(shape, color, *, line=None) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(line)


def _text(tf, text: str, *, size: int, color, bold=False, align=None, font=PPT_FONT) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    _set_run_font(p, name=font)


def _box(slide, l, t, w, h, text: str, *, size, color, bold=False, align=None, wrap=True):
    from pptx.util import Inches

    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    _text(tf, text, size=size, color=color, bold=bold, align=align)
    return box


def _shape(slide, name, l, t, w, h, fill, *, line=None, adj=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    kind = getattr(MSO_SHAPE, name)
    sh = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    _solid(sh, fill, line=line)
    if adj is not None and sh.adjustments:
        try:
            sh.adjustments[0] = adj
        except Exception:
            pass
    return sh


def _bg(slide, width, height, t: dict[str, Any]) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    _solid(sh, t["bg"])
    spPr = sh._element.spPr
    for child in list(spPr):
        if child.tag == qn("a:effectLst"):
            spPr.remove(child)


def _footer(slide, *, deck: str, index: int, total: int, width_in: float, t: dict[str, Any]) -> None:
    _box(slide, 0.7, 7.08, 8.5, 0.28, deck[:40], size=10, color=t["muted"])
    _box(slide, 11.4, 7.08, 1.2, 0.28, f"{index} / {total}", size=10, color=t["muted"], align="right")
    _shape(slide, "RECTANGLE", 0.7, 7.0, width_in - 1.4, 0.015, t["line"])


def _style_cell(cell, text: str, *, fill, fg, bold=False, size=11, center=False) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Pt

    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(fill)
    cell.margin_left = Emu(80000)
    cell.margin_right = Emu(80000)
    cell.margin_top = Emu(50000)
    cell.margin_bottom = Emu(50000)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = _rgb(fg)
        if center:
            p.alignment = PP_ALIGN.CENTER
        _set_run_font(p)


def _add_table(slide, headers, rows, *, left, top, width, height, t: dict[str, Any]):
    from pptx.util import Inches

    ncols = max(len(headers), max((len(r) for r in rows), default=1), 1)
    nrows = 1 + max(len(rows), 1)
    graphic = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = graphic.table
    for j in range(ncols):
        _style_cell(
            tbl.cell(0, j),
            headers[j] if j < len(headers) else "",
            fill=t["accent"],
            fg=t["header_fg"],
            bold=True,
            size=11,
            center=j > 0,
        )
    for i, row in enumerate(rows[: nrows - 1], 1):
        bg = t["surface"] if i % 2 else t["soft"]
        for j in range(ncols):
            val = row[j] if j < len(row) else ""
            _style_cell(
                tbl.cell(i, j),
                str(val),
                fill=bg,
                fg=t["title"] if j == 0 else t["body"],
                bold=j == 0,
                size=12,
                center=j > 0,
            )
    return tbl


def _add_clustered_chart(slide, headers, rows, *, left, top, width, height, t: dict[str, Any]) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt

    names = [str(r[0]) for r in rows if r]
    series_names = [str(h) for h in headers[1:4]]
    if not names or not series_names:
        return
    data = CategoryChartData()
    data.categories = names
    for col, name in enumerate(series_names, 1):
        vals = []
        for r in rows:
            try:
                vals.append(float(r[col]))
            except (IndexError, TypeError, ValueError):
                vals.append(0.0)
        data.add_series(name, vals)
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        data,
    )
    chart = frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)
    plot = chart.plots[0]
    plot.gap_width = 80
    for i, series in enumerate(chart.series):
        color = t["series"][i % len(t["series"])]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _rgb(color)
        series.format.line.fill.background()


def _cover(slide, sl, plan, *, width, height, index, total, t) -> None:
    from pptx.util import Inches, Emu

    _bg(slide, width, height, t)
    _shape(slide, "RECTANGLE", 0, 0, 13.333, 0.1, t["accent"])
    _shape(slide, "RECTANGLE", 0, 0, 0.14, 7.5, t["accent"])
    title = str(sl.get("title") or plan.get("title") or "演示文稿")
    sub = str(sl.get("subtitle") or "")
    points = [str(p).strip() for p in (sl.get("key_points") or []) if str(p).strip()]
    _box(slide, 0.85, 0.55, 6.5, 0.35, "EVALUATION REPORT", size=12, color=t["accent"], bold=True)
    chip = _shape(slide, "ROUNDED_RECTANGLE", 9.3, 0.48, 3.2, 0.42, t["surface"], line=t["line"], adj=0.2)
    _box(slide, 9.3, 0.52, 3.2, 0.36, (points[0] if points else "对比分析")[:22], size=11, color=t["body"], align="center")
    _box(slide, 0.85, 2.15, 11.5, 1.8, title, size=36, color=t["title"], bold=True)
    if sub:
        _box(slide, 0.85, 4.15, 11.2, 0.7, sub, size=18, color=t["body"])
    x = 0.85
    for pill in (points[:4] or ["横向对比"]):
        w = min(3.4, max(1.8, 0.22 * len(pill) + 0.9))
        _shape(slide, "ROUNDED_RECTANGLE", x, 5.15, w, 0.42, t["soft"], line=t["line"], adj=0.5)
        _box(slide, x, 5.2, w, 0.34, pill[:18], size=11, color=t["title"], align="center")
        x += w + 0.18
    _footer(slide, deck=str(plan.get("title") or ""), index=index, total=total, width_in=13.333, t=t)
    _ = (chip, Emu)


def _kpi_cards(slide, sl, *, points, t) -> None:
    n = min(4, max(1, len(points)))
    gap = 0.22
    total_w = 11.9
    w = (total_w - gap * (n - 1)) / n
    y, h = 1.85, 4.55 if n <= 3 else 3.6
    for i, pt in enumerate(points[:n]):
        x = 0.7 + i * (w + gap)
        _shape(slide, "ROUNDED_RECTANGLE", x, y, w, h, t["surface"], line=t["line"], adj=0.06)
        _shape(slide, "RECTANGLE", x, y, 0.08, h, t["accent"] if i % 2 == 0 else t["accent2"])
        head, rest = pt, ""
        if "：" in pt:
            head, rest = pt.split("：", 1)
        elif ":" in pt:
            head, rest = pt.split(":", 1)
        _box(slide, x + 0.22, y + 0.22, w - 0.4, 0.45, f"{i + 1:02d}", size=12, color=t["accent"], bold=True)
        _box(slide, x + 0.22, y + 0.7, w - 0.4, 0.9, head[:36], size=16, color=t["title"], bold=True)
        _box(slide, x + 0.22, y + 1.7, w - 0.4, h - 2.05, (rest or pt)[:120], size=13, color=t["body"])


def _bullet_cards(slide, sl, *, points, t) -> None:
    items = points[:6] or ["按说明补充要点"]
    cols = 3 if len(items) >= 3 else 2 if len(items) == 2 else 1
    rows = (len(items) + cols - 1) // cols
    gap = 0.2
    w = (11.9 - gap * (cols - 1)) / cols
    h = min(2.2, (4.7 - gap * (rows - 1)) / rows)
    for i, pt in enumerate(items):
        r, c = divmod(i, cols)
        x = 0.7 + c * (w + gap)
        y = 1.85 + r * (h + gap)
        _shape(slide, "ROUNDED_RECTANGLE", x, y, w, h, t["surface"], line=t["line"], adj=0.08)
        if "：" in pt:
            head, rest = pt.split("：", 1)
        elif "|" in pt or "｜" in pt:
            parts = [x.strip() for x in re.split(r"[|｜]", pt) if x.strip()]
            head, rest = (parts[0], "  ".join(parts[1:])) if parts else (pt, "")
        else:
            head, rest = pt, ""
        _box(slide, x + 0.2, y + 0.16, w - 0.36, 0.42, f"{i + 1:02d}  {head[:28]}", size=13, color=t["title"], bold=True)
        _box(slide, x + 0.2, y + 0.58, w - 0.36, h - 0.74, (rest or pt)[:100], size=12, color=t["body"])


def _content(slide, sl, plan, *, width, height, index, total, t) -> None:
    _bg(slide, width, height, t)
    title = str(sl.get("title") or "")
    _box(slide, 0.7, 0.38, 11.8, 0.7, title, size=26, color=t["title"], bold=True)
    _shape(slide, "RECTANGLE", 0.7, 1.12, 1.15, 0.07, t["accent"])
    points = [str(p).strip() for p in (sl.get("key_points") or []) if str(p).strip()]
    table = sl.get("table") if isinstance(sl.get("table"), dict) else None
    stype = str(sl.get("type") or "content")
    if not (table and (table.get("headers") or table.get("rows"))):
        pipe_rows = []
        for p in points:
            parts = [x.strip() for x in re.split(r"[|｜]", p) if x.strip()]
            if len(parts) >= 3:
                pipe_rows.append(parts)
        if len(pipe_rows) >= 2:
            width = max(len(r) for r in pipe_rows)
            table = {
                "headers": ["项目", "模型", "得分", "说明"][:width],
                "rows": [r + [""] * (width - len(r)) for r in pipe_rows],
            }
            points = [p for p in points if "|" not in p and "｜" not in p]
            sl = {**sl, "type": "table", "table": table, "key_points": points}

    if table and (table.get("headers") or table.get("rows")):
        headers = [str(h) for h in (table.get("headers") or [])]
        rows = [[str(c) for c in (row if isinstance(row, (list, tuple)) else [row])] for row in (table.get("rows") or [])]
        numeric = any(re.search(r"^0\.\d+", str(c)) for row in rows for c in row[1:])
        if numeric and len(headers) >= 3 and len(rows) >= 2:
            _add_table(slide, headers, rows, left=0.7, top=1.4, width=6.35, height=min(4.9, 0.42 * (len(rows) + 1) + 0.5), t=t)
            _add_clustered_chart(slide, headers, rows, left=7.2, top=1.35, width=5.4, height=5.05, t=t)
        else:
            _add_table(slide, headers, rows, left=0.7, top=1.4, width=11.9, height=min(5.2, 0.45 * (len(rows) + 1) + 0.6), t=t)
    elif stype == "toc" or "目录" in title:
        _bullet_cards(slide, sl, points=points or [str(s.get("title") or "") for s in (plan.get("slides") or [])[1:7]], t=t)
    elif 2 <= len(points) <= 4:
        _kpi_cards(slide, sl, points=points, t=t)
    else:
        _bullet_cards(slide, sl, points=points, t=t)
    _footer(slide, deck=str(plan.get("title") or ""), index=index, total=total, width_in=13.333, t=t)


def build_editable_pptx(plan: dict[str, Any], prs) -> None:
    from pptx.util import Inches

    width, height = Inches(13.333), Inches(7.5)
    prs.slide_width = width
    prs.slide_height = height
    blank = prs.slide_layouts[6]
    t = tokens(str(plan.get("style") or "corporate-clean"))
    slides = list(plan.get("slides") or [])
    total = len(slides)
    for i, sl in enumerate(slides[:24], 1):
        slide = prs.slides.add_slide(blank)
        stype = str(sl.get("type") or "content")
        if stype in {"title"} or i == 1 and stype not in {"content", "table", "toc"}:
            _cover(slide, sl, plan, width=width, height=height, index=i, total=total, t=t)
        else:
            _content(slide, sl, plan, width=width, height=height, index=i, total=total, t=t)
